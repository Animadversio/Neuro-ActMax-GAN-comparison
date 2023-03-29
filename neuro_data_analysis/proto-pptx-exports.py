import os.path
from os.path import join
from pptx.util import Inches
from pptx import Presentation#, SlidePart

#%%
attr_root = r"E:\Network_Data_Sync\BigGAN_FeatAttribution"
protosumdir = r"E:\OneDrive - Harvard University\Manuscript_BigGAN\Figures\ProtoSummary"
#%%
# plt.savefig(join(protosumdir, f"Exp{Expi}_proto_attr_summary.png"))
#%%
ppt_file = Presentation()
# add slide
# slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[6])
#%%
from tqdm import tqdm
for Expi in tqdm(range(1, 190+1)):  # 66 is not good
    imgpath = join(protosumdir, f"Exp{Expi}_proto_attr_summary.png")
    if os.path.exists(imgpath):
        # add slide
        slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[6])
        # add image
        slide.shapes.add_picture(imgpath, left=Inches(0.5), top=Inches(0.0), height=Inches(7.5)) #width=Inches(10),

ppt_file.save(join(protosumdir, "proto_attr_summary.pptx"))