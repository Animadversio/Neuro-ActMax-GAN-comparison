import os.path
import numpy as np
import pandas as pd
from os.path import join
from tqdm import tqdm
from scipy.stats import sem
from pptx.util import Inches
from pptx import Presentation  # , SlidePart
#%%
attr_root = r"E:\Network_Data_Sync\BigGAN_FeatAttribution"
protosumdir = r"E:\OneDrive - Harvard University\Manuscript_BigGAN\Figures\ProtoSummary"
tabdir = r"E:\OneDrive - Harvard University\Manuscript_BigGAN\Stats_tables"
meta_df = pd.read_csv(join(tabdir, "meta_stats.csv"), index_col=False)
meta_act_df = pd.read_csv(join(tabdir, "meta_activation_stats.csv"), index_col=False)
#%% Create masks for summarizing the prototype images
#%%
ppt_file = Presentation()
for Expi in tqdm(range(1, 190+1)):  # 66 is not good
    imgpath = join(protosumdir, f"Exp{Expi}_proto_attr_summary.png")
    if os.path.exists(imgpath):
        # add slide
        slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[6])
        # add image
        slide.shapes.add_picture(imgpath, left=Inches(0.5), top=Inches(0.0), height=Inches(7.5)) #width=Inches(10),

ppt_file.save(join(protosumdir, "proto_attr_summary.pptx"))

#%%
def create_pptx_from_exps(expids):
    ppt_file = Presentation()
    for Expi in tqdm(expids):  # 66 is not good
        imgpath = join(protosumdir, f"Exp{Expi}_proto_attr_summary.png")
        if os.path.exists(imgpath):
            # add slide
            slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[6])
            # add image
            slide.shapes.add_picture(imgpath, left=Inches(0.5), top=Inches(0.0), height=Inches(7.5)) #width=Inches(10),
        else:
            print(f"Exp{Expi} image not found")
    return ppt_file
#%%
succmsk = (meta_act_df.p_maxinit_0 < 0.01) & \
          (meta_act_df.p_maxinit_1 < 0.01)

cmpmsk = (meta_act_df.maxrsp_1_mean - meta_act_df.maxrsp_0_mean).abs() \
       < (meta_act_df.maxrsp_0_sem + meta_act_df.maxrsp_1_sem)

#%%
Expids = meta_act_df[succmsk & cmpmsk].Expi.values
pptx_cmp_bothsucc = create_pptx_from_exps(Expids)
pptx_cmp_bothsucc.save(join(protosumdir, "proto_comparable_success.pptx"))
#%%
pptx_bothsucc = create_pptx_from_exps(meta_act_df[succmsk].Expi.values)
pptx_bothsucc.save(join(protosumdir, "proto_both_success.pptx"))

#%%
imgcmp_root = r"E:\OneDrive - Harvard University\Manuscript_BigGAN\Figures\ProtoImage_cmp"
imgcmpdir = r"E:\OneDrive - Harvard University\Manuscript_BigGAN\Figures\ProtoImage_cmp\scratch"
def create_pptx_from_exps_plus_cmp(expids):
    ppt_file = Presentation()
    for Expi in tqdm(expids):  # 66 is not good
        # if os.path.exists(imgpath):
        # add slide
        slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[6])
        # add image
        slide.shapes.add_picture(join(protosumdir, f"Exp{Expi}_proto_attr_summary.png"),
                                 left=Inches(0.5), top=Inches(0.0), height=Inches(7.5)) #width=Inches(10),

        slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[6])
        slide.shapes.add_picture(join(imgcmpdir, f"Exp{Expi:02d}_FC_BG_reevol_pix.png"),
                                 left=Inches(0.5), top=Inches(0.0), height=Inches(7.5))  # width=Inches(10),

        slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[6])
        slide.shapes.add_picture(join(imgcmpdir, f"Exp{Expi:02d}_FC_BG_reevol_G.png"),
                                 left=Inches(0.5), top=Inches(0.0), height=Inches(7.5))  # width=Inches(10),

        slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[6])
        slide.shapes.add_picture(join(imgcmpdir, f"Exp{Expi:02d}_FC_BG_maxblk.png"),
                                 left=Inches(0.5), top=Inches(0.0), height=Inches(7.5))  # width=Inches(10),
        # else:
        #     print(f"Exp{Expi} image not found")
    return ppt_file


pptx_file = create_pptx_from_exps_plus_cmp(meta_df.Expi.values)
pptx_file.save(join(imgcmp_root, "proto_attr_summary_plus_cmp.pptx"))

#%%
